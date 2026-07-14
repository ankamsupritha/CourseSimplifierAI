import os
import fitz
from docx import Document
from pptx import Presentation
from typing import Any


def extract_text(filepath: str) -> str:

    ext = os.path.splitext(filepath)[1].lower()

    # ---------------- PDF ----------------

    if ext == ".pdf":

        text = ""

        with fitz.open(filepath) as pdf:

            for page in pdf:

                text += str(page.get_text())

        return text

    # ---------------- DOCX ----------------

    elif ext == ".docx":

        document = Document(filepath)

        paragraphs = []

        for para in document.paragraphs:

            if para.text.strip():

                paragraphs.append(para.text)

        return "\n".join(paragraphs)

    # ---------------- PPTX ----------------

    elif ext == ".pptx":

        presentation = Presentation(filepath)

        slides_text = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                text = getattr(shape, "text", None)

                if isinstance(text, str):

                    if text.strip():

                        slides_text.append(text)

        return "\n".join(slides_text)

    # ---------------- TXT ----------------

    elif ext == ".txt":

        with open(filepath, "r", encoding="utf-8") as file:

            return file.read()

    else:

        raise ValueError(f"Unsupported file format: {ext}")